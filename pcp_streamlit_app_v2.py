
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import streamlit as st
from matplotlib.ticker import FuncFormatter
from matplotlib.dates import AutoDateLocator, ConciseDateFormatter
from matplotlib import cycler
from io import BytesIO
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

st.set_page_config(page_title="PCP – Dashboard Dinâmico v6", layout="wide")

# -----------------------------
# Branding (logo)
# -----------------------------
logo_path = Path("icon128.png")
if not logo_path.exists():
    alt = Path("/mnt/data/icon128.png")
    if alt.exists():
        logo_path = alt

# -----------------------------
# Corporate palette (inspired: navy/blue/cyan/light, accent orange)
# -----------------------------
JML_PALETTE = ["#0B1F3B", "#145DA0", "#2E8BC0", "#B1D4E0", "#F2A200"]
plt.rcParams["axes.prop_cycle"] = cycler(color=JML_PALETTE)
plt.rcParams["axes.grid"] = True
plt.rcParams["grid.alpha"] = 0.35
plt.rcParams["grid.linewidth"] = 0.6
plt.rcParams["axes.titleweight"] = "semibold"

# Header with logo
colA, colB = st.columns([0.1, 0.9])
if logo_path.exists():
    colA.image(str(logo_path), use_container_width=True)
colB.title("PCP – Dashboard Dinâmico v6")
st.caption("Paleta custom + export PPTX + metas dinâmicas por família/célula.")

# -----------------------------
# Helpers
# -----------------------------
def format_brl(x, pos=None, dec=0):
    try:
        s = f"{x:,.{dec}f}"
        s = s.replace(",", "X").replace(".", ",").replace("X", ".")
        return f"R$ {s}"
    except Exception:
        return "R$ 0"

def date_axes(ax):
    locator = AutoDateLocator()
    ax.xaxis.set_major_locator(locator)
    ax.xaxis.set_major_formatter(ConciseDateFormatter(locator))

def annotate_peaks(ax, x, y, topn=3, fmt=lambda v: f"{v:.2f}"):
    if len(y) == 0:
        return
    idx = np.argsort(y)[-topn:]
    for i in idx:
        ax.annotate(fmt(y[i]), (x[i], y[i]), xytext=(0, 8), textcoords="offset points", ha="center", fontsize=9)

def rolling_series(series, window):
    if window is None or window <= 1:
        return None
    return series.rolling(window=window, min_periods=max(1, window//2)).mean()

def load_base():
    up = st.sidebar.file_uploader("Substituir base padrão (CSV do template)", type=["csv"])
    if up is not None:
        df = pd.read_csv(up, parse_dates=["month"])
        st.sidebar.success("Base oficial carregada.")
    else:
        df = pd.read_csv("pcp_data.csv", parse_dates=["month"])
    return df

def fig_to_png_bytes(fig):
    bio = BytesIO()
    fig.savefig(bio, format="png", dpi=200, bbox_inches="tight")
    bio.seek(0)
    return bio

def build_pptx(periodo_txt, kpi_dict, figs, logo_path=None):
    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[1]

    # Title slide
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = "PCP – Relatório Executivo"
    sub = slide.placeholders[1]
    sub.text = f"Período: {periodo_txt}\nExportado via Dashboard Dinâmico v6"
    if logo_path and Path(logo_path).exists():
        slide.shapes.add_picture(str(logo_path), Inches(9), Inches(0.5), width=Inches(1.2))

    # KPI slide
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "KPIs do período"
    tb = slide.shapes.placeholders[1].text_frame
    tb.clear()
    for k, v in kpi_dict.items():
        p = tb.add_paragraph()
        p.text = f"• {k}: {v}"
        p.level = 0

    # Chart slides
    for title, fig_bytes in figs:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = title
        slide.shapes.add_picture(fig_bytes, Inches(0.5), Inches(1.2), width=Inches(9))

    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out

# -----------------------------
# Data & Filters
# -----------------------------
df = load_base().sort_values("month").reset_index(drop=True)

st.sidebar.header("Filtros")
min_d, max_d = df["month"].min().date(), df["month"].max().date()
preset = st.sidebar.selectbox("Preset de período", ["Custom", "Últimos 6 meses", "Últimos 12 meses"])
if preset == "Últimos 6 meses":
    start = pd.to_datetime(max_d) - pd.DateOffset(months=6)
    end = pd.to_datetime(max_d)
elif preset == "Últimos 12 meses":
    start = pd.to_datetime(max_d) - pd.DateOffset(months=12)
    end = pd.to_datetime(max_d)
else:
    start = pd.to_datetime(st.sidebar.date_input("Início", min_d))
    end = pd.to_datetime(st.sidebar.date_input("Fim", max_d))

mask = (df["month"] >= start) & (df["month"] <= end)
dff = df.loc[mask].copy()

# Visual controls
st.sidebar.header("Controles de gráfico")
mv_win = st.sidebar.number_input("Média móvel (meses)", min_value=1, value=3, step=1)
topn = st.sidebar.slider("Anotar Top-N picos", min_value=0, max_value=10, value=3, step=1)

# Tabs
tab1, tab2, tab3, tab4, tab5, tab6 = st.tabs([
    "Overview",
    "Custo de Perda vs. Meta",
    "Simulador — No-Cut <7d",
    "Paradas por Código",
    "Metas Dinâmicas por Subgrupo",
    "Exportar PPTX"
])

# -----------------------------
# Tab 1 — Overview
# -----------------------------
with tab1:
    if dff.empty:
        st.warning("Período sem dados.")
    else:
        lead_time = dff["lead_time_mean"].mean()
        efetividade = dff["efetividade_media"].mean()
        perda_prem_r = dff["perda_prem_R"].sum()
        pct_prem = dff["pct_refugo_prem"].replace(0, np.nan).mean()

        c1, c2, c3, c4 = st.columns(4)
        c1.metric("Lead time médio (dias)", f"{lead_time:.2f}")
        c2.metric("Efetividade média", f"{efetividade:.2f}x")
        c3.metric("Perda prematura (R$)", f"{perda_prem_r:,.0f}".replace(",", "."))
        c4.metric("% refugo por corte prematuro", f"{0 if np.isnan(pct_prem) else pct_prem:.1f}%")

        st.subheader("Lead time médio (dias) — série e MM")
        fig1, ax1 = plt.subplots()
        ax1.plot(dff["month"], dff["lead_time_mean"], marker="o")
        ma_lt = rolling_series(dff["lead_time_mean"], mv_win)
        if ma_lt is not None and mv_win > 1:
            ax1.plot(dff["month"], ma_lt, linestyle="--")
        ax1.set_xlabel("Mês"); ax1.set_ylabel("Dias"); ax1.set_title("Lead time médio (dias)")
        date_axes(ax1)
        if topn > 0:
            annotate_peaks(ax1, dff["month"].values, dff["lead_time_mean"].values, topn=topn, fmt=lambda v: f"{v:.1f}")
        st.pyplot(fig1)

        st.subheader("% do refugo por corte prematuro — série e MM")
        fig2, ax2 = plt.subplots()
        ax2.plot(dff["month"], dff["pct_refugo_prem"], marker="o")
        ma_pct = rolling_series(dff["pct_refugo_prem"], mv_win)
        if ma_pct is not None and mv_win > 1:
            ax2.plot(dff["month"], ma_pct, linestyle="--")
        ax2.set_xlabel("Mês"); ax2.set_ylabel("%"); ax2.set_title("% do refugo por corte prematuro")
        date_axes(ax2)
        if topn > 0:
            annotate_peaks(ax2, dff["month"].values, dff["pct_refugo_prem"].values, topn=topn, fmt=lambda v: f"{v:.1f}%")
        st.pyplot(fig2)

# -----------------------------
# Tab 2 — Custo de Perda vs. Meta
# -----------------------------
with tab2:
    meta = st.number_input("Meta mensal (R$)", min_value=0, value=50000, step=5000)
    dff2 = dff.copy()
    dff2["meta"] = meta
    dff2["gap"] = dff2["perda_prem_R"] - dff2["meta"]

    c1, c2 = st.columns(2)
    c1.metric("Perda total no período (R$)", f"{dff2['perda_prem_R'].sum():,.0f}".replace(",", "."))
    c2.metric("Desvio total vs. meta (R$)", f"{dff2['gap'].sum():,.0f}".replace(",", "."))

    st.subheader("Perda por corte prematuro x meta — Mensal")
    fig3, ax3 = plt.subplots()
    ax3.plot(dff2["month"], dff2["perda_prem_R"], marker="o", label="Perda (R$)")
    ax3.plot(dff2["month"], dff2["meta"], linestyle="--", label="Meta (R$)")
    ax3.set_xlabel("Mês"); ax3.set_ylabel("R$"); ax3.set_title("Perda vs. Meta"); ax3.legend()
    date_axes(ax3)
    st.pyplot(fig3)

# -----------------------------
# Tab 3 — Simulador
# -----------------------------
with tab3:
    colA, colB = st.columns(2)
    compliance = colA.slider("Compliance em No-Cut <7d (%)", min_value=0, max_value=100, value=50, step=5)
    meta_alvo = colB.number_input("Meta mensal alvo (R$)", min_value=0, value=50000, step=5000)
    dff3 = dff.copy(); factor = (100 - compliance) / 100.0
    dff3["perda_prem_R_simulada"] = dff3["perda_prem_R"] * factor

    col_a, col_b, col_c = st.columns(3)
    col_a.metric("Perda atual (R$)", f"{dff['perda_prem_R'].sum():,.0f}".replace(",", "."))
    col_b.metric("Perda simulada (R$)", f"{dff3['perda_prem_R_simulada'].sum():,.0f}".replace(",", "."))
    economia = dff["perda_prem_R"].sum() - dff3["perda_prem_R_simulada"].sum()
    col_c.metric("Economia estimada (R$)", f"{economia:,.0f}".replace(",", "."))

    st.subheader("Perda atual vs. simulada")
    fig4, ax4 = plt.subplots()
    ax4.plot(dff["month"], dff["perda_prem_R"], marker="o", label="Atual")
    ax4.plot(dff3["month"], dff3["perda_prem_R_simulada"], linestyle="--", label="Simulada")
    ma_atual = rolling_series(dff["perda_prem_R"], mv_win)
    ma_sim = rolling_series(dff3["perda_prem_R_simulada"], mv_win)
    if ma_atual is not None and mv_win > 1: ax4.plot(dff["month"], ma_atual, linestyle=":")
    if ma_sim is not None and mv_win > 1: ax4.plot(dff3["month"], ma_sim, linestyle=":")
    ax4.set_xlabel("Mês"); ax4.set_ylabel("R$"); ax4.set_title("Perda — Atual vs. Simulada"); ax4.legend()
    date_axes(ax4)
    st.pyplot(fig4)

# -----------------------------
# Tab 4 — Paradas por Código
# -----------------------------
with tab4:
    st.write("Upload de CSV com colunas: **data, codigo_parada, minutos, celula**.")
    up = st.file_uploader("Enviar CSV de Paradas", type=["csv"], key="paradas")
    if up is not None:
        dfp = pd.read_csv(up, parse_dates=["data"])
        msk = (dfp["data"] >= pd.to_datetime(start)) & (dfp["data"] <= pd.to_datetime(end))
        dfp = dfp.loc[msk].copy()

        st.subheader("Top 10 códigos de parada (min)")
        top = dfp.groupby("codigo_parada", as_index=False)["minutos"].sum().sort_values("minutos", ascending=False).head(10)
        st.dataframe(top)

        st.subheader("Série temporal — minutos totais de parada (por mês)")
        dfp["mes"] = dfp["data"].values.astype("datetime64[M]")
        serie = dfp.groupby("mes", as_index=False)["minutos"].sum()

        fig5, ax5 = plt.subplots()
        ax5.plot(serie["mes"], serie["minutos"], marker="o")
        ax5.set_xlabel("Mês"); ax5.set_ylabel("Minutos de parada"); ax5.set_title("Paradas — minutos por mês")
        date_axes(ax5)
        st.pyplot(fig5)

# -----------------------------
# Tab 5 — Metas Dinâmicas por Subgrupo
# -----------------------------
with tab5:
    st.write("Defina metas por **célula** ou **família** via tabela editável e acompanhe gaps e séries.")
    available = [c for c in ["celula","familia"] if c in dff.columns]
    if not available:
        st.info("A base atual não possui 'celula' ou 'familia'. Suba uma base com esses campos na barra lateral.")
    else:
        eixo = st.selectbox("Agrupar por", available)
        unicos = dff[[eixo]].dropna().drop_duplicates().reset_index(drop=True)
        if unicos.empty:
            st.warning("Não há valores no período selecionado para o agrupamento escolhido.")
        else:
            unicos["Meta (R$)"] = 30000
            metas_edit = st.data_editor(unicos, num_rows="dynamic", use_container_width=True, hide_index=True)
            metas_map = dict(zip(metas_edit[eixo], metas_edit["Meta (R$)"]))
            g = dff.groupby([eixo, "month"], as_index=False)["perda_prem_R"].sum()
            g["meta"] = g[eixo].map(metas_map).fillna(0)
            g["gap"] = g["perda_prem_R"] - g["meta"]

            st.subheader("Ranking — maior desvio (gap) no período")
            rank = g.groupby(eixo, as_index=False)["gap"].sum().sort_values("gap", ascending=False)
            rank.columns = [eixo, "Gap total (R$)"]
            st.dataframe(rank, use_container_width=True)

            st.subheader("Série — Perda por subgrupo (linha) e meta (tracejado)")
            fig6, ax6 = plt.subplots()
            for key, chunk in g.groupby(eixo):
                ax6.plot(chunk["month"], chunk["perda_prem_R"], label=str(key))
            for key, chunk in g.groupby(eixo):
                ax6.plot(chunk["month"], chunk["meta"], linestyle="--")
            ax6.set_xlabel("Mês"); ax6.set_ylabel("R$"); ax6.set_title("Perda por subgrupo vs metas"); ax6.legend()
            date_axes(ax6)
            st.pyplot(fig6)

# -----------------------------
# Tab 6 — Exportar PPTX
# -----------------------------
with tab6:
    st.write("Gera um deck executivo (título, KPIs e gráficos do período filtrado).")
    if dff.empty:
        st.info("Selecione um período com dados antes de exportar.")
    else:
        periodo_txt = f"{pd.to_datetime(start).date().strftime('%d/%m/%Y')} — {pd.to_datetime(end).date().strftime('%d/%m/%Y')}"
        # KPIs
        kpi_dict = {
            "Lead time médio (dias)": f"{dff['lead_time_mean'].mean():.2f}",
            "Efetividade média": f"{dff['efetividade_media'].mean():.2f}x",
            "Perda prematura (R$)": f"{dff['perda_prem_R'].sum():,.0f}".replace(",", "."),
            "% refugo por corte prematuro": f"{dff['pct_refugo_prem'].replace(0, np.nan).mean():.1f}%",
        }

        # Replot figs (to ensure clean buffers)
        figs_bytes = []
        # Lead time
        f1, a1 = plt.subplots()
        a1.plot(dff["month"], dff["lead_time_mean"], marker="o")
        ma_lt = rolling_series(dff["lead_time_mean"], mv_win)
        if ma_lt is not None and mv_win > 1: a1.plot(dff["month"], ma_lt, linestyle="--")
        a1.set_xlabel("Mês"); a1.set_ylabel("Dias"); a1.set_title("Lead time médio (dias) — série e MM"); date_axes(a1)
        figs_bytes.append(("Lead time médio (dias)", fig_to_png_bytes(f1)))

        # % refugo prem
        f2, a2 = plt.subplots()
        a2.plot(dff["month"], dff["pct_refugo_prem"], marker="o")
        ma_pct = rolling_series(dff["pct_refugo_prem"], mv_win)
        if ma_pct is not None and mv_win > 1: a2.plot(dff["month"], ma_pct, linestyle="--")
        a2.set_xlabel("Mês"); a2.set_ylabel("%"); a2.set_title("% do refugo por corte prematuro — série e MM"); date_axes(a2)
        figs_bytes.append(("% refugo por corte prematuro", fig_to_png_bytes(f2)))

        # Perda vs meta
        meta = st.number_input("Meta mensal (R$) (para o slide de perda vs meta)", min_value=0, value=50000, step=5000, key="meta_export")
        f3, a3 = plt.subplots()
        a3.plot(dff["month"], dff["perda_prem_R"], marker="o", label="Perda (R$)")
        a3.plot(dff["month"], [meta]*len(dff), linestyle="--", label="Meta (R$)")
        a3.set_xlabel("Mês"); a3.set_ylabel("R$"); a3.set_title("Perda por corte prematuro vs meta"); a3.legend(); date_axes(a3)
        figs_bytes.append(("Perda por corte prematuro vs meta", fig_to_png_bytes(f3)))

        if st.button("Gerar PPTX"):
            pptx_bytes = build_pptx(periodo_txt, kpi_dict, figs_bytes, logo_path=logo_path if logo_path.exists() else None)
            st.download_button("Baixar Relatório Executivo (PPTX)", data=pptx_bytes, file_name="PCP_Relatorio_Executivo.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
